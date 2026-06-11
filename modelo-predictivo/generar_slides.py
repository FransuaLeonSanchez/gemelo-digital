"""Genera la presentación del modelo predictivo de GEMA (4 slides).

Salidas:
    presentacion/img/pipeline_modelo.png   diagrama de arquitectura (slide 1)
    presentacion/slides.html               deck HTML auto-contenido (imágenes embebidas)
    presentacion/GEMA_modelo_ML.pptx       misma presentación en PowerPoint

Requiere los PNG de generar_graficos.py (córrelo antes si no existen).

Uso:
    .venv/bin/python generar_slides.py
"""

import base64
from pathlib import Path

import cairosvg

AQUI = Path(__file__).parent
PRES = AQUI / "presentacion"
IMG = PRES / "img"
IMG.mkdir(parents=True, exist_ok=True)

BG = "#0B0E14"; CARD = "#131A2C"; LINE = "#3A4358"; TXT = "#EAF0FA"; SUB = "#8A95AC"
GREEN = "#1FD0A3"; BLUE = "#4DA3FF"; PURPLE = "#A78BFA"; AMBER = "#FFB23E"; RED = "#FF5E6C"

# ═════════════════════════════════════════════════════════════════════════════
# 1 · Diagrama de pipeline (SVG → PNG)
# ═════════════════════════════════════════════════════════════════════════════

def esc(s): return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def txt(x, y, s, size=13, fill=TXT, weight="600", anchor="middle"):
    return (f'<text x="{x}" y="{y}" font-size="{size}" fill="{fill}" font-weight="{weight}" '
            f'text-anchor="{anchor}" font-family="DejaVu Sans, sans-serif">{esc(s)}</text>')

def card(cx, cy, w, h, color, lines):
    out = [f'<rect x="{cx-w/2}" y="{cy-h/2}" width="{w}" height="{h}" rx="12" '
           f'fill="{CARD}" stroke="{color}" stroke-width="2"/>']
    lh = 19
    y0 = cy - (len(lines) - 1) * lh / 2 + 4
    for i, (s, fill, weight, size) in enumerate(lines):
        out.append(txt(cx, y0 + i * lh, s, size=size, fill=fill, weight=weight))
    return "".join(out)

def flecha(pts, color=SUB):
    p = " ".join(f"{a},{b}" for a, b in pts)
    return (f'<polyline points="{p}" fill="none" stroke="{color}" stroke-width="1.7" '
            f'marker-end="url(#fl)"/>')

W, H = 1560, 640
S = [f'<svg xmlns="http://www.w3.org/2000/svg" width="{W}" height="{H}" viewBox="0 0 {W} {H}">'
     f'<defs><marker id="fl" viewBox="0 0 10 10" refX="9" refY="5" markerWidth="7" '
     f'markerHeight="7" orient="auto-start-reverse">'
     f'<path d="M 0 0 L 10 5 L 0 10 z" fill="{SUB}"/></marker></defs>'
     f'<rect width="{W}" height="{H}" fill="{BG}"/>']

for cx, titulo in [(210, "DE DÓNDE SALEN LOS DATOS"), (590, "QUÉ MIRA EL MODELO"),
                   (975, "LOS 2 MODELOS"), (1365, "QUÉ ENTREGA EL MODELO")]:
    S.append(txt(cx, 56, titulo, size=14, fill=SUB, weight="800"))

S.append(card(210, 110, 310, 76, PURPLE, [
    ("Foto del plato (IA de la app)", PURPLE, "800", 13.5),
    ("estima carbohidratos, índice", SUB, "600", 12),
    ("glucémico y carga del plato", SUB, "600", 12)]))
S.append(card(210, 206, 310, 76, BLUE, [
    ("Smartwatch", BLUE, "800", 13.5),
    ("pasos · horas de sueño · caminata", SUB, "600", 12),
    ("variabilidad del ritmo cardíaco", SUB, "600", 12)]))
S.append(card(210, 302, 310, 76, AMBER, [
    ("Sensor de glucosa (CGM)", AMBER, "800", 13.5),
    ("glucosa basal · tiempo en valores", SUB, "600", 12),
    ("saludables · variabilidad", SUB, "600", 12)]))
S.append(card(210, 398, 310, 76, GREEN, [
    ("Formulario de la app", GREEN, "800", 13.5),
    ("medida de cintura", SUB, "600", 12)]))
S.append(card(210, 494, 310, 76, RED, [
    ("Análisis de sangre (opcional)", RED, "800", 13.5),
    ("triglicéridos y glucosa en ayunas", SUB, "600", 12),
    ("→ índice TyG", SUB, "600", 12)]))

S.append(card(590, 200, 290, 86, LINE, [
    ("RESUMEN DE 7 DÍAS", TXT, "800", 13.5),
    ("8 datos del usuario:", SUB, "600", 12),
    ("hábitos + glucosa + cuerpo", SUB, "600", 12)]))
S.append(card(590, 400, 290, 86, LINE, [
    ("LA COMIDA QUE VA A COMER", TXT, "800", 13.5),
    ("8 datos del plato", SUB, "600", 12),
    ("y del momento", SUB, "600", 12)]))

S.append(card(975, 200, 320, 96, GREEN, [
    ("MODELO 1 · RANDOM FOREST", GREEN, "800", 13.5),
    ("nivel de riesgo: bajo · moderado · alto", TXT, "600", 12.5),
    ("acierta 9 de cada 10 (88.5 %)", SUB, "600", 12)]))
S.append(card(975, 400, 320, 96, BLUE, [
    ("MODELO 2 · GRADIENT BOOSTING", BLUE, "800", 13.5),
    ("pico de glucosa a los 60 min", TXT, "600", 12.5),
    ("error promedio ±6 mg/dL", SUB, "600", 12)]))

S.append(card(1365, 200, 300, 84, GREEN, [
    ("Nivel de riesgo + probabilidad", GREEN, "800", 13),
    ("bajo · moderado · alto", TXT, "600", 12),
    ("→ ICM proyectado y gemelo", SUB, "600", 12)]))
S.append(card(1365, 400, 300, 84, BLUE, [
    ("Pico de glucosa a 60 min", BLUE, "800", 13),
    ("en mg/dL, al registrar el plato", TXT, "600", 12),
    ("→ se refleja en el ICM del día", SUB, "600", 12)]))

# fuentes → features (canales verticales propios para evitar superposición)
S.append(flecha([(365, 190), (383, 190), (383, 170), (442, 170)]))   # watch → resumen
S.append(flecha([(365, 286), (391, 286), (391, 190), (442, 190)]))   # sensor → resumen
S.append(flecha([(365, 398), (399, 398), (399, 210), (442, 210)]))   # formulario → resumen
S.append(flecha([(365, 494), (407, 494), (407, 230), (442, 230)]))   # análisis → resumen
S.append(flecha([(365, 110), (417, 110), (417, 376), (442, 376)]))   # foto → comida
S.append(flecha([(365, 222), (425, 222), (425, 400), (442, 400)]))   # watch → comida
S.append(flecha([(365, 318), (433, 318), (433, 424), (442, 424)]))   # sensor → comida
# features → modelos
S.append(flecha([(735, 200), (812, 200)]))
S.append(flecha([(735, 400), (812, 400)]))
# modelos → salidas
S.append(flecha([(1135, 200), (1212, 200)]))
S.append(flecha([(1135, 400), (1212, 400)]))

S.append(txt(W / 2, 596, "Circuito completo de punta a punta: cuando lleguen datos de usuarios reales, solo cambia la fuente de datos — el motor ya está listo.",
             size=13, fill=SUB, weight="600"))
S.append("</svg>")

cairosvg.svg2png(bytestring="".join(S).encode(), write_to=str(IMG / "pipeline_modelo.png"), scale=2)
print("✔ pipeline_modelo.png")

# ═════════════════════════════════════════════════════════════════════════════
# 2 · Contenido compartido por HTML y PPTX
# ═════════════════════════════════════════════════════════════════════════════

SLIDE2_COLS = [
    ("QUÉ RECIBE · 8 datos de los últimos 7 días", BLUE, [
        "Sensor de glucosa: tiempo en valores saludables, variabilidad y pico promedio después de comer",
        "Smartwatch: horas de sueño, pasos diarios y variabilidad del ritmo cardíaco (indicador de estrés)",
        "Análisis de sangre (opcional): índice TyG — triglicéridos y glucosa en ayunas",
        "Formulario de la app: medida de cintura",
    ]),
    ("QUÉ DEVUELVE", GREEN, [
        "Nivel de riesgo: bajo / moderado / alto",
        "La probabilidad de cada nivel (qué tan seguro está el modelo)",
        "Con esto la app calcula el ICM proyectado: el color y la expresión del gemelo",
    ]),
    ("QUÉ TAN BIEN FUNCIONA · 1 200 casos de prueba", AMBER, [
        "Acierta el nivel de riesgo en 9 de cada 10 evaluaciones (88.5 %)",
        "El resultado es estable al repetir la prueba 5 veces (87.5 ± 0.9 %)",
        "Nunca confunde riesgo bajo con alto: si falla, es entre niveles vecinos",
    ]),
]

SLIDE3_COLS = [
    ("QUÉ RECIBE · 8 datos de la comida y el momento", BLUE, [
        "Del plato (foto analizada por la IA de la app): carbohidratos, índice glucémico y carga del plato",
        "Del momento: glucosa al sentarse a comer y hora del día",
        "Del smartwatch: caminata después de comer, sueño de anoche y ritmo cardíaco",
    ]),
    ("QUÉ DEVUELVE", GREEN, [
        "El pico de glucosa que ese plato causará ~60 minutos después, en mg/dL",
        "Con esto la app actualiza el ICM del día y el estado del gemelo",
    ]),
    ("QUÉ TAN BIEN FUNCIONA · 1 800 comidas de prueba", AMBER, [
        "Error promedio de ±6 mg/dL — la meta clínica era ±15",
        "9 de cada 10 predicciones (93.8 %) caen dentro del margen aceptado de ±15 mg/dL",
    ]),
]

SLIDE4_CHIPS = [
    ("Hoy", GREEN,
     "Entrenado con población simulada, construida con estudios poblacionales (Perú y EE. UU.) y fisiología publicada"),
    ("Piloto", BLUE,
     "2 semanas de calibración por usuario, con sensor de glucosa real y smartwatch"),
    ("Futuro", PURPLE,
     "El modelo aprende continuamente de cada usuario y afina sus predicciones"),
]

SLIDE1_CHIPS = [
    ("Con qué se entrenó", GREEN,
     "6 000 perfiles y 9 000 comidas simuladas a partir de estudios poblacionales (Perú y EE. UU.) y fisiología publicada, mientras llegan los datos de usuarios reales"),
    ("Por qué es confiable", BLUE,
     "Resultados verificables y repetibles: el sistema completo se puede volver a correr y entrega siempre las mismas cifras"),
    ("Metas cumplidas", AMBER,
     "Riesgo: acierta 88.5 % (meta: 85 %) · Glucosa: error de ±6 mg/dL (meta: ±15) — ambas superadas"),
]

TITULOS = {
    1: ("GEMA · El cerebro del gemelo digital",
        "Dos modelos de aprendizaje automático convierten tus datos diarios en predicciones útiles"),
    2: ("Modelo 1 · ¿Qué nivel de riesgo metabólico tienes hoy?",
        "Random Forest: combina 200 árboles de decisión para clasificar el riesgo con los últimos 7 días del usuario"),
    3: ("Modelo 2 · ¿Cuánto te subirá la glucosa este plato?",
        "Gradient Boosting: anticipa el pico de glucosa antes de comer, para decidir a tiempo"),
    4: ("El modelo en acción",
        "Las predicciones se convierten en decisiones concretas y medibles — y así sigue creciendo"),
}
N_SLIDES = 4

# ═════════════════════════════════════════════════════════════════════════════
# 3 · Deck HTML auto-contenido
# ═════════════════════════════════════════════════════════════════════════════

def b64(name):
    return "data:image/png;base64," + base64.b64encode((IMG / name).read_bytes()).decode()

def html_cols(cols):
    out = ['<div class="cols">']
    for titulo, color, items in cols:
        lis = "".join(f"<li>{i}</li>" for i in items)
        out.append(f'<div class="col" style="border-top:3px solid {color}">'
                   f'<h3 style="color:{color}">{titulo}</h3><ul>{lis}</ul></div>')
    out.append("</div>")
    return "".join(out)

def html_header(n, kicker_color=GREEN):
    t, s = TITULOS[n]
    return (f'<header><div><h1>{t}</h1><p class="sub">{s}</p></div>'
            f'<span class="num" style="color:{kicker_color}">{n} / {N_SLIDES}</span></header>')

html = f"""<!DOCTYPE html>
<html lang="es"><head><meta charset="utf-8">
<title>GEMA · Modelo predictivo del gemelo digital</title>
<style>
  * {{ margin:0; padding:0; box-sizing:border-box; }}
  body {{ background:#05070B; font-family:'Segoe UI','DejaVu Sans',sans-serif;
         display:flex; flex-direction:column; align-items:center; gap:28px; padding:28px; }}
  .slide {{ width:1280px; height:720px; background:{BG}; border:1px solid {LINE};
           border-radius:18px; padding:30px 40px 26px; display:flex; flex-direction:column;
           overflow:hidden; position:relative; }}
  header {{ display:flex; justify-content:space-between; align-items:flex-start; }}
  h1 {{ color:{TXT}; font-size:27px; font-weight:800; letter-spacing:-0.3px; }}
  .sub {{ color:{SUB}; font-size:15px; margin-top:5px; }}
  .num {{ font-size:13px; font-weight:800; border:1px solid {LINE}; border-radius:99px;
         padding:5px 12px; background:{CARD}; white-space:nowrap; }}
  .cols {{ display:flex; gap:20px; margin-top:16px; }}
  .col {{ flex:1; background:{CARD}; border:1px solid {LINE}; border-radius:14px;
         padding:14px 16px; }}
  .col h3 {{ font-size:13px; font-weight:800; letter-spacing:0.6px; margin-bottom:8px; }}
  .col ul {{ list-style:none; }}
  .col li {{ color:{TXT}; font-size:13.5px; line-height:1.42; margin-bottom:7px;
            padding-left:14px; position:relative; }}
  .col li::before {{ content:"•"; color:{SUB}; position:absolute; left:0; }}
  .charts {{ flex:1; display:flex; gap:24px; justify-content:center; align-items:center;
            margin-top:10px; min-height:0; }}
  .charts img {{ max-height:100%; max-width:48%; border-radius:12px; }}
  .pipeline {{ flex:1; display:flex; justify-content:center; align-items:center; min-height:0; }}
  .pipeline img {{ max-width:100%; max-height:100%; }}
  .chips {{ display:flex; gap:18px; }}
  .chip {{ flex:1; background:{CARD}; border:1px solid {LINE}; border-radius:14px;
          padding:12px 16px; }}
  .chip b {{ display:block; font-size:12px; letter-spacing:0.6px; margin-bottom:5px;
            text-transform:uppercase; }}
  .chip span {{ color:{TXT}; font-size:13px; line-height:1.4; }}
  footer {{ position:absolute; bottom:10px; left:40px; right:40px; display:flex;
           justify-content:space-between; color:{SUB}; font-size:11px; font-weight:600; }}
</style></head><body>

<section class="slide" id="s1">
  {html_header(1)}
  <div class="pipeline"><img src="{b64('pipeline_modelo.png')}" alt="pipeline"></div>
  <div class="chips">
    {"".join(f'<div class="chip" style="border-top:3px solid {c}"><b style="color:{c}">{t}</b><span>{s}</span></div>' for t, c, s in SLIDE1_CHIPS)}
  </div>
  <footer><span>GEMA · Gemelo digital metabólico — Lima, Perú</span><span>Junio 2026</span></footer>
</section>

<section class="slide" id="s2">
  {html_header(2)}
  {html_cols(SLIDE2_COLS)}
  <div class="charts">
    <img src="{b64('confusion_riesgo.png')}" alt="aciertos y errores" style="max-width:56%">
  </div>
  <footer><span>GEMA · Gemelo digital metabólico — Lima, Perú</span><span>Modelo 1 · Random Forest</span></footer>
</section>

<section class="slide" id="s3">
  {html_header(3)}
  {html_cols(SLIDE3_COLS)}
  <div class="charts">
    <img src="{b64('pred_vs_real_pico.png')}" alt="predicción vs realidad" style="max-width:56%">
  </div>
  <footer><span>GEMA · Gemelo digital metabólico — Lima, Perú</span><span>Modelo 2 · Gradient Boosting</span></footer>
</section>

<section class="slide" id="s4">
  {html_header(4)}
  <div class="charts">
    <img src="{b64('escenarios_whatif.png')}" alt="el mismo almuerzo, cuatro decisiones">
    <img src="{b64('importancias_riesgo.png')}" alt="qué datos pesan más">
  </div>
  <div class="chips">
    {"".join(f'<div class="chip" style="border-top:3px solid {c}"><b style="color:{c}">{t}</b><span>{s}</span></div>' for t, c, s in SLIDE4_CHIPS)}
  </div>
  <footer><span>GEMA · Gemelo digital metabólico — Lima, Perú</span><span>Del dato a la decisión</span></footer>
</section>

</body></html>"""

(PRES / "slides.html").write_text(html, encoding="utf-8")
print("✔ slides.html")

# ═════════════════════════════════════════════════════════════════════════════
# 4 · PPTX (mismo contenido, editable en PowerPoint)
# ═════════════════════════════════════════════════════════════════════════════

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

def C(hexs): return RGBColor.from_string(hexs.lstrip("#"))
PX = lambda p: Emu(int(p / 96 * 914400))   # 1280×720 px → 13.33×7.5 in

prs = Presentation()
prs.slide_width = PX(1280)
prs.slide_height = PX(720)
blank = prs.slide_layouts[6]

def nueva_slide():
    s = prs.slides.add_slide(blank)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = C(BG)
    return s

def caja_texto(slide, x, y, w, h):
    tb = slide.shapes.add_textbox(PX(x), PX(y), PX(w), PX(h))
    tb.text_frame.word_wrap = True
    return tb.text_frame

def parrafo(tf, texto, size, color, bold=False, first=False, space_after=4):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    r = p.add_run(); r.text = texto
    f = r.font; f.size = Pt(size); f.color.rgb = C(color); f.bold = bold
    f.name = "Segoe UI"
    p.space_after = Pt(space_after)
    return p

def tarjeta(slide, x, y, w, h, accent):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PX(x), PX(y), PX(w), PX(h))
    sh.adjustments[0] = 0.07
    sh.fill.solid(); sh.fill.fore_color.rgb = C(CARD)
    sh.line.color.rgb = C(accent); sh.line.width = Pt(1.4)
    sh.shadow.inherit = False
    return sh

def cabecera(slide, n):
    t, s = TITULOS[n]
    tf = caja_texto(slide, 40, 22, 1110, 78)
    parrafo(tf, t, 24, TXT, bold=True, first=True, space_after=2)
    parrafo(tf, s, 12.5, SUB)
    tf2 = caja_texto(slide, 1160, 30, 90, 30)
    p = parrafo(tf2, f"{n} / {N_SLIDES}", 11, GREEN, bold=True, first=True)
    p.alignment = PP_ALIGN.RIGHT

def pie(slide, izq, der):
    tf = caja_texto(slide, 40, 692, 800, 24)
    parrafo(tf, izq, 9, SUB, first=True)
    tf2 = caja_texto(slide, 880, 692, 360, 24)
    p = parrafo(tf2, der, 9, SUB, first=True)
    p.alignment = PP_ALIGN.RIGHT

def columnas(slide, cols, y=112, h=190):
    xs = [40, 456, 872]; w = 392
    for (titulo, color, items), x in zip(cols, xs):
        tarjeta(slide, x, y, w, h, color)
        tf = caja_texto(slide, x + 14, y + 10, w - 28, h - 20)
        parrafo(tf, titulo, 11.5, color, bold=True, first=True, space_after=6)
        for it in items:
            parrafo(tf, "• " + it, 11, TXT, space_after=4)

def imagenes(slide, nombres, y=318, alto=360):
    from PIL import Image
    anchos = []
    for n in nombres:
        with Image.open(IMG / n) as im:
            anchos.append(alto * im.width / im.height)
    gap = 30
    total = sum(anchos) + gap * (len(nombres) - 1)
    x = (1280 - total) / 2
    for n, w in zip(nombres, anchos):
        slide.shapes.add_picture(str(IMG / n), PX(x), PX(y), height=PX(alto))
        x += w + gap

# Slide 1
s1 = nueva_slide()
cabecera(s1, 1)
from PIL import Image
with Image.open(IMG / "pipeline_modelo.png") as im:
    pw = 1130; ph = pw * im.height / im.width
s1.shapes.add_picture(str(IMG / "pipeline_modelo.png"), PX((1280 - pw) / 2), PX(108), width=PX(pw))
y_chips = 108 + ph + 14
xs = [40, 456, 872]; w = 392; hch = 700 - y_chips - 24
for (titulo, color, texto), x in zip(SLIDE1_CHIPS, xs):
    tarjeta(s1, x, y_chips, w, hch, color)
    tf = caja_texto(s1, x + 14, y_chips + 8, w - 28, hch - 16)
    parrafo(tf, titulo.upper(), 10.5, color, bold=True, first=True, space_after=4)
    parrafo(tf, texto, 10.5, TXT)
pie(s1, "GEMA · Gemelo digital metabólico — Lima, Perú", "Junio 2026")

# Slide 2
s2 = nueva_slide()
cabecera(s2, 2)
columnas(s2, SLIDE2_COLS)
imagenes(s2, ["confusion_riesgo.png"])
pie(s2, "GEMA · Gemelo digital metabólico — Lima, Perú", "Modelo 1 · Random Forest")

# Slide 3
s3 = nueva_slide()
cabecera(s3, 3)
columnas(s3, SLIDE3_COLS)
imagenes(s3, ["pred_vs_real_pico.png"])
pie(s3, "GEMA · Gemelo digital metabólico — Lima, Perú", "Modelo 2 · Gradient Boosting")

# Slide 4
s4 = nueva_slide()
cabecera(s4, 4)
imagenes(s4, ["escenarios_whatif.png", "importancias_riesgo.png"], y=112, alto=330)
y_road = 470
for (titulo, color, texto), x in zip(SLIDE4_CHIPS, [40, 456, 872]):
    tarjeta(s4, x, y_road, 392, 130, color)
    tf = caja_texto(s4, x + 14, y_road + 10, 392 - 28, 130 - 20)
    parrafo(tf, titulo.upper(), 11, color, bold=True, first=True, space_after=5)
    parrafo(tf, texto, 11, TXT)
pie(s4, "GEMA · Gemelo digital metabólico — Lima, Perú", "Del dato a la decisión")

prs.save(PRES / "GEMA_modelo_ML.pptx")
print("✔ GEMA_modelo_ML.pptx")
